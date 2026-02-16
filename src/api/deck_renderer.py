"""
Klaus x IMI — Deck Renderer (python-pptx)
==========================================
Core rendering engine. Takes structured slide configs (JSON dicts)
and produces IMI-branded .pptx slides with charts, stat cards,
narrative headlines, and SO WHAT callouts.

Usage:
    from deck_renderer import build_deck
    build_deck(slide_configs, "output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData

# ═══════════════════════════════════════════════════════════════
# IMI BRAND CONSTANTS
# ═══════════════════════════════════════════════════════════════

NAVY       = RGBColor(0x0B, 0x1D, 0x3A)
DARK_NAVY  = RGBColor(0x07, 0x14, 0x28)
CARD_NAVY  = RGBColor(0x0F, 0x25, 0x49)
ORANGE     = RGBColor(0xE8, 0x65, 0x1A)
ORANGE_LT  = RGBColor(0xF2, 0x8C, 0x4E)
TEAL       = RGBColor(0x00, 0xA3, 0xA1)
GOLD       = RGBColor(0xD4, 0xA8, 0x43)
GREEN      = RGBColor(0x2E, 0x8B, 0x57)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF5, 0xF6, 0xFA)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF2)
MED_GRAY   = RGBColor(0x8C, 0x99, 0xA9)
DARK_GRAY  = RGBColor(0x3A, 0x4A, 0x5C)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

FONT_H = "Georgia"
FONT_B = "Calibri"

ACCENTS = [TEAL, ORANGE, GOLD, GREEN]


# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
    s.fill.solid()
    s.fill.fore_color.rgb = ORANGE
    s.line.fill.background()

def _txt(slide, l, t, w, h, text, font=FONT_B, sz=13, color=DARK_GRAY,
         bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.name = font
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    return tb

def _rich(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for i, r in enumerate(runs):
        run = p.runs[0] if i == 0 and p.runs else p.add_run()
        run.text = r.get("text", "")
        run.font.name = r.get("font", FONT_B)
        run.font.size = Pt(r.get("sz", 11))
        run.font.color.rgb = r.get("color", DARK_GRAY)
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
    return tb

def _card(slide, l, t, w, h, accent, bg=WHITE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t + 0.05), Inches(0.07), Inches(h - 0.1))
    b.fill.solid()
    b.fill.fore_color.rgb = accent
    b.line.fill.background()
    return c

def _source(slide, text, dark=False):
    c = MED_GRAY if not dark else RGBColor(0x5A, 0x6A, 0x7A)
    _txt(slide, 0.7, 5.2, 8.6, 0.3, text, sz=9, color=c)

def _insight(slide, text, y=4.65, dark=False):
    label_c = ORANGE
    text_c = WHITE if dark else DARK_GRAY
    _rich(slide, 0.7, y, 8.6, 0.5, [
        {"text": "SO WHAT? ", "bold": True, "color": label_c, "sz": 10},
        {"text": text, "color": text_c, "sz": 10}
    ])


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════

def slide_title(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY); _bar(s)
    _txt(s, 0.7, 0.3, 5, 0.3, "INSIGHT. DRIVING. PROFIT.",
         sz=12, color=ORANGE, bold=True)
    _txt(s, 0.7, 1.4, 7, 1.2, c.get("title", "IMI Pulse™ Analysis"),
         font=FONT_H, sz=44, color=WHITE, bold=True)
    _txt(s, 0.7, 2.7, 7, 0.6, c.get("subtitle", ""),
         sz=16, color=ORANGE_LT, italic=True)
    _txt(s, 0.7, 4.8, 8.6, 0.4, c.get("methodology", ""),
         sz=11, color=MED_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════

def slide_exec_summary(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, OFF_WHITE); _bar(s)
    _txt(s, 0.7, 0.2, 8, 0.5, "EXECUTIVE SUMMARY",
         font=FONT_H, sz=22, color=NAVY, bold=True)
    _txt(s, 0.7, 0.65, 8, 0.3, c.get("exec_subtitle", "The Essential Insight"),
         sz=12, color=ORANGE, italic=True)

    for i, stat in enumerate(c.get("stats", [])[:3]):
        y = 1.2 + i * 1.25
        accent = ACCENTS[i % len(ACCENTS)]
        _card(s, 0.7, y, 8.6, 1.05, accent)
        _txt(s, 1.0, y + 0.1, 1.5, 0.7, stat.get("number", ""),
             font=FONT_H, sz=36, color=accent, bold=True)
        _txt(s, 2.6, y + 0.15, 6.5, 0.75, stat.get("description", ""),
             sz=13, color=DARK_GRAY)

    _source(s, c.get("source", "Source: IMI Pulse™"))


# ═══════════════════════════════════════════════════════════════
# BAR CHART SLIDE (used for slides 3, 4, 8)
# ═══════════════════════════════════════════════════════════════

def slide_bar_chart(prs, c, bar_color=ORANGE, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, bg); _bar(s)
    _txt(s, 0.7, 0.2, 8.6, 0.8, c.get("title", ""),
         font=FONT_H, sz=20, color=NAVY, bold=True)
    _txt(s, 0.7, 0.85, 8.6, 0.3, c.get("subtitle", ""),
         sz=11, color=MED_GRAY, italic=True)

    data = c.get("data", {})
    if data:
        cd = CategoryChartData()
        labels = list(reversed(list(data.keys())))
        vals = list(reversed(list(data.values())))
        cd.categories = labels
        cd.add_series('', vals)

        cf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.7), Inches(1.25), Inches(8.6), Inches(3.3), cd)
        ch = cf.chart
        ch.has_legend = False

        plot = ch.plots[0]
        plot.gap_width = 80
        ser = plot.series[0]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = bar_color

        ser.has_data_labels = True
        dl = ser.data_labels
        dl.font.size = Pt(11)
        dl.font.bold = True
        dl.font.color.rgb = NAVY
        dl.number_format = '0"%"'
        dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

        ch.category_axis.tick_labels.font.size = Pt(10)
        ch.category_axis.tick_labels.font.color.rgb = DARK_GRAY
        ch.category_axis.has_major_gridlines = False

        va = ch.value_axis
        va.visible = False
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = LIGHT_GRAY
        va.major_gridlines.format.line.width = Pt(0.5)

    ins = c.get("insight", "")
    if ins:
        _rich(s, 0.7, 4.65, 8.6, 0.5, [
            {"text": "KEY INSIGHT: ", "bold": True, "color": ORANGE, "sz": 10},
            {"text": ins, "color": DARK_GRAY, "sz": 10}
        ])
    _source(s, c.get("source", "Source: IMI Pulse™"))


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: NATIVITY DIVIDE (dark, side-by-side tables)
# ═══════════════════════════════════════════════════════════════

def slide_nativity(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY); _bar(s)
    _txt(s, 0.7, 0.2, 8.6, 0.5, c.get("title", "THE NATIVITY DIVIDE"),
         font=FONT_H, sz=24, color=WHITE, bold=True)
    _txt(s, 0.7, 0.65, 8.6, 0.3, c.get("subtitle", ""),
         sz=12, color=ORANGE_LT, italic=True)

    def _col(x_start, label, n, data_dict, accent):
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x_start), Inches(1.1), Inches(4.1), Inches(0.35))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = accent
        hdr.line.fill.background()
        _txt(s, x_start + 0.1, 1.12, 3.9, 0.3,
             f"{label} (n={n})", sz=11, color=WHITE, bold=True)

        y = 1.55
        for i, (lbl, val) in enumerate(data_dict.items()):
            row_bg = CARD_NAVY if i % 2 == 0 else NAVY
            row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(x_start), Inches(y), Inches(4.1), Inches(0.35))
            row.fill.solid()
            row.fill.fore_color.rgb = row_bg
            row.line.fill.background()

            vc = accent if i == 0 else WHITE
            _txt(s, x_start + 0.1, y + 0.02, 2.8, 0.3, lbl, sz=10, color=WHITE)
            _txt(s, x_start + 3.1, y + 0.02, 0.9, 0.3,
                 f"{val}%", sz=11, color=vc, bold=True, align=PP_ALIGN.RIGHT)
            y += 0.35

    _col(0.7,
         c.get("left_label", "BORN IN CANADA"),
         c.get("left_n", ""),
         c.get("left_data", {}),
         TEAL)
    _col(5.2,
         c.get("right_label", "BORN OUTSIDE CANADA"),
         c.get("right_n", ""),
         c.get("right_data", {}),
         ORANGE)

    so = c.get("so_what", "")
    if so:
        _insight(s, so, dark=True)


# ═══════════════════════════════════════════════════════════════
# GROUPED BAR CHART (used for slides 6, 9)
# ═══════════════════════════════════════════════════════════════

def slide_grouped_bar(prs, c, bg=OFF_WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, bg); _bar(s)
    _txt(s, 0.7, 0.2, 8.6, 0.8, c.get("title", ""),
         font=FONT_H, sz=20, color=NAVY, bold=True)
    _txt(s, 0.7, 0.85, 8.6, 0.3, c.get("subtitle", ""),
         sz=11, color=MED_GRAY, italic=True)

    grouped = c.get("data", {})
    colors_hex = c.get("colors", ["E8651A", "00A3A1", "0B1D3A"])

    if grouped:
        cd = CategoryChartData()
        cats = list(grouped.keys())
        cd.categories = cats

        first_cat = list(grouped.values())[0]
        series_names = list(first_cat.keys())

        for sn in series_names:
            cd.add_series(sn, [grouped[cat].get(sn, 0) for cat in cats])

        cf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.7), Inches(1.25), Inches(8.6), Inches(3.0), cd)
        ch = cf.chart
        ch.has_legend = True
        ch.legend.position = 2
        ch.legend.font.size = Pt(10)
        ch.legend.include_in_layout = False

        plot = ch.plots[0]
        plot.gap_width = 60

        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = RGBColor.from_string(
                colors_hex[i % len(colors_hex)])
            ser.has_data_labels = True
            dl = ser.data_labels
            dl.font.size = Pt(9)
            dl.font.bold = True
            dl.font.color.rgb = NAVY
            dl.number_format = '0"%"'
            dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

        ch.value_axis.visible = False
        ch.value_axis.has_major_gridlines = True
        ch.value_axis.major_gridlines.format.line.color.rgb = LIGHT_GRAY
        ch.category_axis.tick_labels.font.size = Pt(10)
        ch.category_axis.tick_labels.font.color.rgb = DARK_GRAY

    ins = c.get("insight", "")
    if ins:
        _rich(s, 0.7, 4.45, 8.6, 0.5, [
            {"text": "KEY INSIGHT: ", "bold": True, "color": ORANGE, "sz": 10},
            {"text": ins, "color": DARK_GRAY, "sz": 10}
        ])
    _source(s, c.get("source", "Source: IMI Pulse™"))


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: GENDER & REGIONAL
# ═══════════════════════════════════════════════════════════════

def slide_gender_regional(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _bar(s)
    _txt(s, 0.7, 0.2, 8.6, 0.5, c.get("title", "GENDER & REGIONAL LENS"),
         font=FONT_H, sz=20, color=NAVY, bold=True)

    # Gender table (left)
    gd = c.get("gender_data", [])
    if gd:
        rows = len(gd) + 1
        tbl = s.shapes.add_table(rows, 4,
            Inches(0.7), Inches(1.0), Inches(4.2), Inches(0.35 * rows)).table

        for j, h in enumerate(["Event", "Female", "Male", "Gap"]):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_B; p.font.size = Pt(10)
                p.font.color.rgb = WHITE; p.font.bold = True
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

        for i, rd in enumerate(gd):
            for j, val in enumerate(rd):
                cell = tbl.cell(i + 1, j)
                cell.text = str(val)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFB)
                for p in cell.text_frame.paragraphs:
                    p.font.name = FONT_B; p.font.size = Pt(10)
                    p.font.color.rgb = DARK_GRAY
                    p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

        for j, w in enumerate([Inches(1.8), Inches(0.8), Inches(0.8), Inches(0.8)]):
            tbl.columns[j].width = w

    # Regional cards (right)
    for i, reg in enumerate(c.get("regions", [])[:4]):
        y = 1.0 + i * 0.9
        _card(s, 5.3, y, 4.0, 0.75, ORANGE)
        _txt(s, 5.55, y + 0.08, 3.6, 0.25, reg.get("name", ""),
             sz=11, color=ORANGE, bold=True)
        _txt(s, 5.55, y + 0.35, 3.6, 0.35, reg.get("insight", ""),
             sz=10, color=DARK_GRAY)

    so = c.get("so_what", "")
    if so:
        _insight(s, so)
    _source(s, c.get("source", "Source: IMI Pulse™"))


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: STRATEGIC IMPLICATIONS (dark, 2x2 cards)
# ═══════════════════════════════════════════════════════════════

def slide_strategic(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY); _bar(s)
    _txt(s, 0.7, 0.2, 8.6, 0.5, c.get("title", "STRATEGIC IMPLICATIONS"),
         font=FONT_H, sz=22, color=WHITE, bold=True)
    _txt(s, 0.7, 0.65, 8.6, 0.3, "Turning Insight Into Action",
         sz=12, color=ORANGE_LT, italic=True)

    pos = [(0.7,1.15,4.2,1.7), (5.1,1.15,4.2,1.7),
           (0.7,3.05,4.2,1.7), (5.1,3.05,4.2,1.7)]

    for i, rec in enumerate(c.get("recommendations", [])[:4]):
        x, y, w, h = pos[i]
        accent = ACCENTS[i % len(ACCENTS)]

        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_NAVY
        card.line.fill.background()

        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y + 0.05), Inches(0.07), Inches(h - 0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        _txt(s, x+0.2, y+0.15, w-0.4, 0.4, rec.get("title",""),
             sz=11, color=accent, bold=True)
        _txt(s, x+0.2, y+0.55, w-0.4, h-0.7, rec.get("body",""),
             sz=10, color=WHITE)

    _source(s, c.get("source", "Source: IMI Pulse™"), dark=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: CLOSING
# ═══════════════════════════════════════════════════════════════

def slide_closing(prs, c):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK_NAVY); _bar(s)
    _txt(s, 0.7, 1.5, 8.6, 0.5, "INSIGHT. DRIVING. PROFIT.",
         sz=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 1.5, 2.2, 7, 0.7,
         c.get("cta_title", "Let's Unlock the Value of Your Data"),
         font=FONT_H, sz=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 1.5, 3.0, 7, 0.5,
         c.get("cta_subtitle", "For a deeper dive, reach out to our team."),
         sz=13, color=MED_GRAY, align=PP_ALIGN.CENTER)

    btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.8), Inches(3.8), Inches(2.4), Inches(0.5))
    btn.fill.solid()
    btn.fill.fore_color.rgb = ORANGE
    btn.line.fill.background()
    tf = btn.text_frame
    p = tf.paragraphs[0]
    p.text = "consultimi.com"
    p.font.name = FONT_B; p.font.size = Pt(14)
    p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    _txt(s, 1.5, 4.8, 7, 0.3,
         "IMI International | Toronto • Melbourne • Tokyo",
         sz=10, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# MAIN BUILD
# ═══════════════════════════════════════════════════════════════

def build_deck(configs: dict, output_path: str = "output.pptx") -> str:
    """
    Build complete 11-slide IMI deck from slide configs dict.

    configs keys: slide_1 through slide_11
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, configs.get("slide_1", {}))
    slide_exec_summary(prs, configs.get("slide_2", {}))
    slide_bar_chart(prs, configs.get("slide_3", {}), bar_color=ORANGE)
    slide_bar_chart(prs, configs.get("slide_4", {}), bar_color=NAVY)
    slide_nativity(prs, configs.get("slide_5", {}))
    slide_grouped_bar(prs, configs.get("slide_6", {}))
    slide_gender_regional(prs, configs.get("slide_7", {}))
    slide_bar_chart(prs, configs.get("slide_8", {}), bar_color=TEAL)
    slide_grouped_bar(prs, configs.get("slide_9", {}), bg=OFF_WHITE)
    slide_strategic(prs, configs.get("slide_10", {}))
    slide_closing(prs, configs.get("slide_11", {}))

    prs.save(output_path)
    return output_path
