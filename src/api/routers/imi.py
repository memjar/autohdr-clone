"""
IMI Router — All /klaus/imi/* endpoints.
Extracted from main.py for clean separation of concerns.

13 endpoints:
  POST /klaus/imi/chat              — Main chat (RAG + pipeline + streaming)
  POST /klaus/imi/report            — Full report generation
  POST /klaus/imi/visualize         — Chart config JSON
  GET  /klaus/imi/chart-types       — Available chart types
  GET  /klaus/imi/dashboard         — Pre-computed brand metrics
  GET  /klaus/imi/stats             — System health
  POST /klaus/imi/ingest            — RAG re-ingestion
  POST /klaus/imi/generate-training — QLoRA training pair generation
  POST /klaus/imi/save              — Save conversation
  POST /klaus/imi/upload-survey     — Upload survey file
  POST /klaus/imi/generate-deck     — Generate PPTX from survey
  POST /klaus/imi/demo-deck         — Generate PPTX from cached analysis
  GET  /klaus/imi/surveys           — List loaded surveys
"""

import csv
import io
import json
import os
import logging
import re

from datetime import datetime
from fastapi import APIRouter, File, UploadFile, HTTPException, Request
from fastapi.responses import StreamingResponse, Response, JSONResponse

from config import (
    OLLAMA_BASE, IMI_MODEL, IMI_DATA_DIR, VECTOR_STORE_PATH,
    LEARNING_LOG_PATH, TRAINING_FILE_PATH, IMI_MEMORY_DIR,
    IMI_NUM_CTX, IMI_NUM_PREDICT, TOOL_CALLING_ENABLED,
)
from ollama_client import async_client
from agent_loop import agent_loop_streaming

logger = logging.getLogger("imi")

router = APIRouter(prefix="/klaus/imi", tags=["imi"])

# In-memory survey store (shared state — imported by main.py)
SURVEY_STORE: dict = {}


# ─── Helpers ──────────────────────────────────────────────


def _load_team_context(max_messages: int = 30) -> str:
    """Load recent team channel messages for Klaus memory context."""
    from config import TEAM_CHANNEL_PATH, ACTIVE_CONTEXT_PATH
    lines = []
    try:
        with open(TEAM_CHANNEL_PATH, "r") as f:
            all_lines = f.readlines()
            for raw in all_lines[-max_messages:]:
                try:
                    m = json.loads(raw.strip())
                    lines.append(f"[{m.get('from','?')}→{m.get('to','team')}] {m.get('msg','')[:200]}")
                except Exception:
                    pass
    except Exception:
        pass

    context_parts = []
    if lines:
        context_parts.append("Recent team channel (Forge, Cortana, Klaus, James):\n" + "\n".join(lines))
    try:
        with open(ACTIVE_CONTEXT_PATH, "r") as f:
            ctx = json.load(f)
            if ctx.get("current_projects"):
                context_parts.append("Active projects: " + json.dumps(ctx["current_projects"][:5]))
            if ctx.get("critical_facts"):
                context_parts.append("Critical facts: " + json.dumps(ctx["critical_facts"][:5]))
    except Exception:
        pass
    return "\n\n".join(context_parts)


def _stream_text(text: str, model: str = IMI_MODEL, chunk_size: int = 3):
    """Convert a string into Ollama-compatible NDJSON stream."""
    import asyncio

    async def _stream():
        words = text.split(" ")
        for i in range(0, len(words), chunk_size):
            chunk = " ".join(words[i : i + chunk_size])
            if i > 0:
                chunk = " " + chunk
            yield json.dumps({"model": model, "message": {"role": "assistant", "content": chunk}, "done": False}).encode() + b"\n"
            await asyncio.sleep(0.02)
        yield json.dumps({"model": model, "message": {"role": "assistant", "content": ""}, "done": True}).encode() + b"\n"

    return _stream()


# ─── Survey Parsers ───────────────────────────────────────


def _parse_crosstab_xlsx(file_bytes: bytes) -> dict:
    """Parse IMI crosstab xlsx (segments as columns, Q-blocks as rows)."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    ws = wb.active
    rows = [[cell.value for cell in row] for row in ws.iter_rows()]
    if len(rows) < 4:
        return {"total_n": 0, "segments": [], "questions": [], "format": "crosstab"}

    segments = [str(c).strip() for c in rows[0] if c and str(c).strip() and str(c).strip().lower() not in ("none", "")]
    questions = []
    current_q = None

    for r in rows[3:]:
        label = str(r[0]).strip() if r[0] else ""
        if not label or label.lower() in ("none", "nan"):
            continue
        if label.startswith("Q") and any(c is not None for c in r[1:]):
            current_q = {"id": label.split(".")[0] if "." in label else label, "text": label, "options": [], "base_sizes": {}}
            questions.append(current_q)
        elif current_q is not None:
            vals = {}
            for j, seg in enumerate(segments):
                idx = j + 1
                if idx < len(r) and r[idx] is not None:
                    try:
                        vals[seg] = float(r[idx])
                    except (ValueError, TypeError):
                        vals[seg] = str(r[idx])
            if vals:
                current_q["options"].append({"label": label, "values": vals})

    return {"total_n": 0, "segments": segments, "questions": questions, "format": "crosstab"}


def _parse_flat_xlsx(file_bytes: bytes) -> dict:
    """Parse flat xlsx — headers in row 1, data in subsequent rows."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    ws = wb.active
    rows_data = [[cell.value for cell in row] for row in ws.iter_rows()]
    if len(rows_data) < 2:
        return {"total_n": 0, "segments": ["Total"], "questions": [], "format": "flat_xlsx"}

    headers = [str(h).strip() for h in rows_data[0] if h]
    questions = []
    for h in headers:
        opts = []
        for row in rows_data[1:]:
            idx = rows_data[0].index(h) if h in [str(x) for x in rows_data[0]] else -1
            if idx >= 0 and idx < len(row) and row[idx] is not None:
                opts.append({"label": str(row[idx]), "values": {"Total": 1}})
        if opts:
            questions.append({"id": h, "text": h.replace("_", " ").title(), "options": opts[:20], "base_sizes": {}})

    return {"total_n": len(rows_data) - 1, "segments": ["Total"], "questions": questions, "format": "flat_xlsx"}


def _parse_flat_csv(file_bytes: bytes) -> dict:
    """Parse flat CSV file."""
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.DictReader(io.StringIO(text))
    headers = reader.fieldnames or []
    rows = list(reader)

    questions = []
    for h in headers:
        vals = {}
        for row in rows:
            v = row.get(h, "")
            if v:
                vals[v] = vals.get(v, 0) + 1
        top = sorted(vals.items(), key=lambda x: -x[1])[:20]
        if top:
            questions.append({
                "id": h,
                "text": h.replace("_", " ").title(),
                "options": [{"label": k, "values": {"Total": v}} for k, v in top],
                "base_sizes": {},
            })

    return {"total_n": len(rows), "segments": ["Total"], "questions": questions, "format": "flat_csv"}


def _parse_json_survey(file_bytes: bytes) -> dict:
    """Parse JSON survey data."""
    data = json.loads(file_bytes.decode("utf-8"))
    items = []
    questions = []

    if isinstance(data, list):
        items = data
        if items and isinstance(items[0], dict):
            for key in list(items[0].keys())[:30]:
                vals = {}
                for item in items:
                    v = str(item.get(key, ""))
                    if v and v != "None":
                        vals[v] = vals.get(v, 0) + 1
                top = sorted(vals.items(), key=lambda x: -x[1])[:20]
                if top:
                    questions.append({
                        "id": key,
                        "text": key.replace("_", " ").title(),
                        "options": [{"label": k, "values": {"Total": v}} for k, v in top],
                        "base_sizes": {},
                    })
    elif isinstance(data, dict):
        items = list(data.values()) if all(isinstance(v, dict) for v in data.values()) else [data]
        for mk, mv in data.items():
            opts = []
            if isinstance(mv, dict):
                for sk, sv in list(mv.items())[:20]:
                    opts.append({"label": sk, "values": {"Total": sv if isinstance(sv, (int, float)) else 1}})
            if opts:
                questions.append({
                    "id": mk,
                    "text": mk.replace(".", " > ").replace("_", " ").title(),
                    "options": opts,
                    "base_sizes": {},
                })

    return {"total_n": 0, "segments": ["Total"], "questions": questions, "format": "nested_json", "items": items}


# ─── Endpoints ────────────────────────────────────────────


@router.post("/save")
async def save_conversation(request: Request):
    """Auto-save IMI conversation to memory."""
    body = await request.json()
    conv_id = body.get("conversation_id", "unknown")
    title = body.get("title", "untitled")
    messages = body.get("messages", [])
    filepath = os.path.join(IMI_MEMORY_DIR, f"{conv_id}.json")
    with open(filepath, "w") as f:
        json.dump({"id": conv_id, "title": title, "messages": messages, "saved_at": datetime.now().isoformat()}, f, indent=2)
    return {"status": "saved", "path": filepath}


@router.get("/surveys")
async def list_surveys():
    """List loaded surveys."""
    return {
        "surveys": [
            {
                "id": sid,
                "name": s.get("structured_data", {}).get("survey_name", sid),
                "total_n": s.get("structured_data", {}).get("total_n", 0),
            }
            for sid, s in SURVEY_STORE.items()
        ]
    }


@router.post("/upload-survey")
async def upload_survey(file: UploadFile):
    """Upload survey file (xlsx/csv/json), parse it, return structured data."""
    fname = file.filename or ""
    if not fname.endswith((".xlsx", ".xls", ".csv", ".json")):
        raise HTTPException(status_code=400, detail="Must be .xlsx, .csv, or .json")

    file_bytes = await file.read()
    survey_name = fname.rsplit(".", 1)[0].replace("_", " ").strip()
    if "book" in survey_name.lower():
        survey_name = "Canadian Sports Fandom Pulse"

    try:
        if fname.endswith((".xlsx", ".xls")):
            parsed = _parse_crosstab_xlsx(file_bytes)
            if not parsed["questions"]:
                parsed = _parse_flat_xlsx(file_bytes)
        elif fname.endswith(".csv"):
            parsed = _parse_flat_csv(file_bytes)
        elif fname.endswith(".json"):
            parsed = _parse_json_survey(file_bytes)
        else:
            raise HTTPException(status_code=400, detail="Unsupported format")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Parse error: {str(e)}")

    survey_data = {
        "survey_name": survey_name,
        "total_n": parsed.get("total_n", 0),
        "segments": parsed.get("segments", []),
        "questions": parsed.get("questions", []),
        "format": parsed.get("format", "unknown"),
        "items": parsed.get("items", []),
    }
    sid = survey_name.lower().replace(" ", "_")
    SURVEY_STORE[sid] = {"structured_data": survey_data, "raw_text": json.dumps(survey_data, default=str)}

    return {
        "survey_id": sid,
        "survey_name": survey_name,
        "total_n": survey_data["total_n"],
        "questions_found": len(survey_data["questions"]),
        "segments": survey_data["segments"],
        "format": survey_data["format"],
    }


@router.post("/generate-deck")
async def generate_deck(request: Request):
    """Generate an IMI-branded PPTX deck from a loaded survey."""
    body = await request.json()
    survey_id = body.get("survey_id")
    if not survey_id or survey_id not in SURVEY_STORE:
        raise HTTPException(status_code=404, detail=f"Survey '{survey_id}' not loaded. Available: {list(SURVEY_STORE.keys())}")

    survey = SURVEY_STORE[survey_id]
    survey_data = survey.get("structured_data")
    if not survey_data:
        raise HTTPException(status_code=400, detail="No structured data for this survey")

    from imi_deck_generator import generate_imi_deck

    pptx_bytes = generate_imi_deck(survey_data)
    fname = survey_data.get("survey_name", "IMI_Analysis").replace(" ", "_") + "_IMI.pptx"
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'},
    )


@router.post("/demo-deck")
async def demo_deck(request: Request):
    """Generate PPTX deck from cached demo analysis markdown."""
    body = await request.json()
    prompt = body.get("prompt", "")
    dataset_name = body.get("dataset_name", "IMI Analysis")

    try:
        from imi_cached_responses import CACHED_RESPONSES
    except ImportError:
        return JSONResponse({"error": "Cached responses not available"}, status_code=500)

    if prompt not in CACHED_RESPONSES:
        return JSONResponse({"error": "No cached response for this prompt"}, status_code=404)

    analysis_md = CACHED_RESPONSES[prompt]

    from imi_deck_generator import (
        _add_shape, _add_text_shape, _orange_bar, _source_line, _insight_box,
        NAVY, ORANGE, TEAL, WHITE, GREY, LIGHT_BG, DARK_FOOTER,
        SLIDE_W, SLIDE_H, MARGIN, CONTENT_W,
    )
    from pptx import Presentation
    from pptx.util import Pt, Emu

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Parse markdown
    title_match = re.search(r"^#\s+.*?—\s*(.+)", analysis_md, re.MULTILINE)
    title = title_match.group(1).strip() if title_match else dataset_name

    exec_match = re.search(r"## Executive Summary\s*\n\n(.+?)(?=\n\n---|\n\n##)", analysis_md, re.DOTALL)
    exec_summary = exec_match.group(1).strip() if exec_match else ""

    sowhat_match = re.search(r"\*\*SO WHAT\?\*\*\s*(.+?)(?=\n\n###|\n\n\*Source)", analysis_md, re.DOTALL)
    sowhat = sowhat_match.group(1).strip() if sowhat_match else ""

    source_match = re.search(r"\*Source:(.+?)\*", analysis_md, re.DOTALL)
    source = "Source:" + source_match.group(1).strip() if source_match else "Source: IMI Pulse™"

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(457200), CONTENT_W, Emu(365760), "INSIGHT. DRIVING. PROFIT.", font_name="Georgia", font_size=Pt(11), color=ORANGE, bold=True)
    _add_text_shape(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(914400), title, font_name="Georgia", font_size=Pt(36), bold=True, color=WHITE)
    _add_text_shape(slide, MARGIN, Emu(2514600), CONTENT_W, Emu(365760), f"IMI Pulse™ Analysis  |  {dataset_name}", font_size=Pt(14), color=GREY)
    _add_text_shape(slide, MARGIN, Emu(4114800), CONTENT_W, Emu(274320), "Powered by Klaus — IMI Intelligence Engine", font_size=Pt(10), color=GREY)
    _add_shape(slide, 0, Emu(int(SLIDE_H) - 91440), SLIDE_W, Emu(91440), DARK_FOOTER)
    _add_text_shape(slide, MARGIN, Emu(int(SLIDE_H) - 82296), Emu(4572000), Emu(73152), "consultimi.com", font_size=Pt(8), color=GREY)

    # Executive Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(182880), CONTENT_W, Emu(365760), "EXECUTIVE SUMMARY", font_name="Georgia", font_size=Pt(22), bold=True, color=WHITE)
    _add_text_shape(slide, MARGIN, Emu(640080), CONTENT_W, Emu(3200400), exec_summary.replace("**", "").replace("*", ""), font_size=Pt(12), color=WHITE)
    _source_line(slide, source)
    _add_shape(slide, 0, Emu(int(SLIDE_H) - 91440), SLIDE_W, Emu(91440), DARK_FOOTER)

    # Content slides from ## sections
    section_pattern = r"##\s+(.+?)\n\n(.*?)(?=\n##\s|\n\*Source:|\Z)"
    all_sections = re.findall(section_pattern, analysis_md, re.DOTALL)

    for sec_title, sec_content in all_sections[:8]:
        if sec_title == "Executive Summary":
            continue
        if "Strategic Implications" in sec_title or "Recommended Actions" in sec_title:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
        _orange_bar(slide)
        _add_text_shape(slide, MARGIN, Emu(182880), CONTENT_W, Emu(365760), sec_title.replace("**", "").strip().upper(), font_name="Georgia", font_size=Pt(18), bold=True, color=WHITE)

        table_match = re.search(r"\|(.+)\|\n\|[-\s|:]+\|\n((?:\|.+\|\n)*)", sec_content)
        if table_match:
            headers = [h.strip() for h in table_match.group(1).split("|") if h.strip()]
            rows = []
            for row_line in table_match.group(2).strip().split("\n"):
                cells = [c.strip().replace("**", "") for c in row_line.split("|") if c.strip()]
                if cells:
                    rows.append(cells)

            table_text = "  ".join(h[:15].ljust(15) for h in headers[:6]) + "\n"
            table_text += "─" * min(90, len(headers[:6]) * 17) + "\n"
            for row in rows[:10]:
                table_text += "  ".join(str(c)[:15].ljust(15) for c in row[:6]) + "\n"

            _add_text_shape(slide, MARGIN, Emu(640080), CONTENT_W, Emu(3200400), table_text, font_size=Pt(9), color=WHITE)
        else:
            clean = sec_content.replace("**", "").replace("*", "").replace("> ", "")[:800]
            _add_text_shape(slide, MARGIN, Emu(640080), CONTENT_W, Emu(3200400), clean, font_size=Pt(11), color=WHITE)

        insight_match = re.search(r"\*\*(?:Key (?:Insight|Finding)|Insight):\*\*\s*(.+)", sec_content)
        if insight_match:
            _insight_box(slide, insight_match.group(1).strip()[:200])

        _source_line(slide, source)
        _add_shape(slide, 0, Emu(int(SLIDE_H) - 91440), SLIDE_W, Emu(91440), DARK_FOOTER)

    # Strategic Implications slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(182880), CONTENT_W, Emu(365760), "STRATEGIC IMPLICATIONS", font_name="Georgia", font_size=Pt(22), bold=True, color=WHITE)
    if sowhat:
        _add_text_shape(slide, MARGIN, Emu(640080), CONTENT_W, Emu(1371600), sowhat.replace("**", "").replace("*", "").replace("> ", "")[:600], font_size=Pt(13), color=WHITE)
    actions_match = re.search(r"### Recommended Actions\s*\n((?:\d+\..+\n?)*)", analysis_md)
    if actions_match:
        _add_text_shape(slide, MARGIN, Emu(2194560), CONTENT_W, Emu(1828800), actions_match.group(1).replace("**", "").strip()[:500], font_size=Pt(11), color=ORANGE)
    _add_shape(slide, 0, Emu(int(SLIDE_H) - 91440), SLIDE_W, Emu(91440), DARK_FOOTER)

    # Closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(548640), "THANK YOU", font_name="Georgia", font_size=Pt(36), bold=True, color=WHITE)
    _add_text_shape(slide, MARGIN, Emu(2057400), CONTENT_W, Emu(365760), "For more information, contact IMI International", font_size=Pt(14), color=GREY)
    _add_text_shape(slide, MARGIN, Emu(2514600), CONTENT_W, Emu(274320), "consultimi.com  |  info@consultimi.com", font_size=Pt(12), color=TEAL)
    _add_text_shape(slide, MARGIN, Emu(3474720), CONTENT_W, Emu(274320), "Analysis powered by Klaus — IMI Intelligence Engine", font_size=Pt(10), color=GREY)
    _add_shape(slide, 0, Emu(int(SLIDE_H) - 91440), SLIDE_W, Emu(91440), DARK_FOOTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    filename = dataset_name.replace(" ", "_") + "_IMI_Deck.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/chat")
async def chat(request: Request):
    """Klaus IMI chat — RAG + multi-agent pipeline + Ollama streaming."""
    import asyncio

    body = await request.json()
    message = body.get("message", "")
    history = body.get("history", [])

    # Check cached responses first (zero compute)
    try:
        from imi_cached_responses import CACHED_RESPONSES
        if message in CACHED_RESPONSES and not history:
            return StreamingResponse(_stream_text(CACHED_RESPONSES[message]), status_code=200, media_type="application/x-ndjson")
    except ImportError:
        pass

    # RAG retrieval (capped at 8 chunks to keep context small and fast)
    rag_context = ""
    rag_chunks = []
    try:
        from imi_rag import query as rag_query, format_context
        rag_chunks = rag_query(message, n=8)
        if rag_chunks:
            rag_context = format_context(rag_chunks)
    except Exception as e:
        logger.warning(f"RAG query failed: {e}")

    # Multi-agent pipeline ONLY for truly complex multi-dataset queries
    # Simple questions go straight to Ollama (1 LLM call vs 2+)
    pipeline_triggers = [
        "compare", "comparison", "report", "analyze across", "full analysis",
        "cross-dataset", "deep dive", "gap analysis", "break down",
    ]
    if any(t in message.lower() for t in pipeline_triggers) and rag_chunks:
        try:
            from imi_agents import run_pipeline
            result = run_pipeline(message, rag_chunks)
            if result and len(result) > 50:
                return StreamingResponse(_stream_text(result), status_code=200, media_type="application/x-ndjson")
        except Exception as e:
            logger.warning(f"Pipeline failed, falling back to direct: {e}")

    # Build system message (kept concise to reduce token count)
    system_msg = (
        "You are Klaus, IMI International's AI insight engine (18 countries, 55+ years research).\n"
        "Bold key numbers (**47%**). Cite dataset names. NEVER fabricate data.\n"
        "Format: **Executive Summary** (2-3 sentences) | **Key Data** (table) | **Analysis** | **SO WHAT?** (3 actions)"
    )
    if rag_context:
        system_msg += f"\n\n{rag_context}"

    # Only load team context for meta/team queries — saves ~500+ tokens on every request
    if any(t in message.lower() for t in ["team", "forge", "cortana", "project", "axe"]):
        team_context = _load_team_context(max_messages=10)
        if team_context:
            system_msg += f"\n\n{team_context}"

    file_catalog = body.get("file_catalog", "")
    if file_catalog:
        system_msg += f"\n\nFILE CATALOG (files available in the system):\n{file_catalog}\nYou can reference these files when users ask what data is available."

    survey_id = body.get("survey_id")
    if survey_id and survey_id in SURVEY_STORE:
        survey = SURVEY_STORE[survey_id]
        data_inject = json.dumps(survey.get("structured_data"), indent=2) if survey.get("structured_data") else survey.get("raw_text", "")
        system_msg += f"\n\n<survey_data>\n{data_inject}\n</survey_data>\nThe above is the ONLY data you may reference. Every number in your response must come from this data."

    messages = [{"role": "system", "content": system_msg}]
    messages += [{"role": h.get("role", "user"), "content": h.get("content", "")} for h in history]
    messages.append({"role": "user", "content": message})

    # Temperature based on query type
    temp = 0.1
    if any(t in message.lower() for t in ["deck", "presentation", "slides", "full analysis", "insight report", "analyze"]):
        temp = 0.4
    elif any(t in message.lower() for t in ["should", "recommend", "strategy", "implication", "opportunity"]):
        temp = 0.3

    # Use agent loop with tool calling if enabled
    use_tools = body.get("use_tools", TOOL_CALLING_ENABLED)

    if use_tools:
        # Agent loop: streams NDJSON events (tool_start, tool_result, token, thinking, done)
        user_messages = [{"role": h.get("role", "user"), "content": h.get("content", "")} for h in history]
        user_messages.append({"role": "user", "content": message})

        async def stream_agent():
            async for event in agent_loop_streaming(
                messages=user_messages,
                system_prompt=system_msg,
                model=IMI_MODEL,
                max_iterations=5,
                tools_enabled=True,
            ):
                yield json.dumps(event).encode() + b"\n"

        return StreamingResponse(stream_agent(), status_code=200, media_type="application/x-ndjson")

    # Fallback: direct Ollama streaming (no tools)
    ollama_payload = {
        "model": IMI_MODEL,
        "messages": messages,
        "stream": True,
        "keep_alive": "24h",
        "options": {"temperature": temp, "top_p": 0.9, "num_ctx": IMI_NUM_CTX, "num_predict": IMI_NUM_PREDICT},
    }

    req = async_client.build_request(
        "POST", f"{OLLAMA_BASE}/api/chat",
        content=json.dumps(ollama_payload).encode(),
        headers={"content-type": "application/json"},
    )
    resp = await async_client.send(req, stream=True)

    async def stream_ollama():
        try:
            async for line in resp.aiter_lines():
                yield line.encode() + b"\n"
        finally:
            await resp.aclose()

    return StreamingResponse(stream_ollama(), status_code=resp.status_code, media_type="application/x-ndjson")


@router.post("/ingest")
async def ingest(request: Request):
    """Ingest IMI datasets into vector store for RAG."""
    body = await request.json() if request.headers.get("content-type") == "application/json" else {}
    force = body.get("force", False)
    try:
        from imi_rag import ingest_datasets
        result = ingest_datasets(force=force)
        return {"ok": True, **result}
    except Exception as e:
        return {"ok": False, "error": str(e)}


@router.get("/dashboard")
async def dashboard():
    """Pre-computed brand metrics from IMI CSVs."""
    brands = []
    datasets_info = []
    total_records = 0

    if not os.path.isdir(IMI_DATA_DIR):
        return {"brands": [], "datasets": 0, "total_records": 0, "datasets_info": []}

    for filename in sorted(os.listdir(IMI_DATA_DIR)):
        filepath = os.path.join(IMI_DATA_DIR, filename)
        if filename.endswith(".csv"):
            with open(filepath, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                rows = list(reader)
                total_records += len(rows)
                datasets_info.append({"name": filename, "records": len(rows), "type": "csv"})
                if filename == "brand_health_tracker_Q4_2025.csv":
                    for row in rows:
                        brands.append({
                            "brand": row.get("Brand", ""),
                            "nps": int(row.get("NPS", 0)),
                            "awareness": round(float(row.get("Aided Awareness", 0)) * 100),
                            "loyalty": round(float(row.get("Loyalty Score", 0)) * 100),
                            "category": row.get("Category", ""),
                        })
        elif filename.endswith(".json"):
            try:
                with open(filepath, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    count = len(data) if isinstance(data, list) else 1
                    total_records += count
                    datasets_info.append({"name": filename, "records": count, "type": "json"})
            except Exception:
                datasets_info.append({"name": filename, "records": 0, "type": "json"})

    return {"brands": brands, "datasets": len(datasets_info), "total_records": total_records, "datasets_info": datasets_info}


@router.post("/report")
async def report(request: Request):
    """Generate a structured IMI report using RAG + multi-agent pipeline."""
    body = await request.json()
    topic = body.get("topic", "brand health overview")

    try:
        from imi_rag import query as rag_query
        from imi_agents import run_pipeline

        rag_chunks = rag_query(topic, n=10)
        if not rag_chunks:
            return {"ok": False, "error": "No relevant data found for this topic"}

        report_text = run_pipeline(topic, rag_chunks)
        datasets_used = list(set(c.get("dataset", "") for c in rag_chunks))
        return {"ok": True, "topic": topic, "report": report_text, "datasets_used": datasets_used, "chunks_retrieved": len(rag_chunks)}
    except Exception as e:
        logger.error(f"Report generation failed: {e}")
        return {"ok": False, "error": str(e)}


@router.post("/visualize")
async def visualize(request: Request):
    """Return chart config JSON for frontend rendering."""
    body = await request.json()
    chart_type = body.get("type", "nps_comparison")

    chart_handlers = {
        "nps_comparison": _chart_nps,
        "awareness": _chart_awareness,
        "market_share": _chart_market_share,
        "sponsorship_roi": _chart_sponsorship,
        "sports_viewership": _chart_sports,
        "nativity_divide": _chart_nativity,
        "promotion_roi": _chart_promotion_roi,
        "say_do_gap": _chart_say_do,
    }

    handler = chart_handlers.get(chart_type)
    if handler:
        return handler()

    return {"error": f"Unknown chart type: {chart_type}", "available": list(chart_handlers.keys())}


def _read_csv(filename: str) -> list[dict]:
    """Read a CSV from the IMI data directory."""
    filepath = os.path.join(IMI_DATA_DIR, filename)
    if not os.path.exists(filepath):
        return []
    with open(filepath, "r", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def _chart_nps():
    data = [{"name": r.get("Brand", ""), "nps": int(r.get("NPS", 0))} for r in _read_csv("brand_health_tracker_Q4_2025.csv")]
    data.sort(key=lambda x: x["nps"], reverse=True)
    return {"chart": "bar", "title": "NPS by Brand — Q4 2025", "data": data, "xKey": "name", "yKey": "nps"}


def _chart_awareness():
    data = [
        {"name": r.get("Brand", ""), "aided": round(float(r.get("Aided Awareness", 0)) * 100), "unaided": round(float(r.get("Unaided Awareness", 0)) * 100)}
        for r in _read_csv("brand_health_tracker_Q4_2025.csv")
    ]
    data.sort(key=lambda x: x["aided"], reverse=True)
    return {"chart": "bar", "title": "Brand Awareness — Q4 2025", "data": data, "xKey": "name", "yKeys": ["aided", "unaided"]}


def _chart_market_share():
    data = [
        {"name": r.get("Brand", ""), "value": round(float(r.get("Market_Share_Pct", r.get("Market Share", "0"))), 1)}
        for r in _read_csv("competitive_benchmark_market_share.csv")
    ]
    data.sort(key=lambda x: x["value"], reverse=True)
    return {"chart": "pie", "title": "Market Share", "data": data[:10]}


def _chart_sponsorship():
    data = [
        {"name": r.get("Property", r.get("Brand", "")), "score": round(float(r.get("Opportunity_Score", r.get("ROI_Score", "0"))), 1)}
        for r in _read_csv("sponsorship_property_scores.csv")
    ]
    data.sort(key=lambda x: x["score"], reverse=True)
    return {"chart": "bar", "title": "Sponsorship Property Scores", "data": data[:10], "xKey": "name", "yKey": "score"}


def _chart_sports():
    rows = _read_csv("sports_viewership_crosstab_2025.csv")
    data = [
        {"name": r.get("Response", ""), "value": round(float(r.get("Proportion", "0")) * 100, 1)}
        for r in rows
        if r.get("Question", "").startswith("Q36") and r.get("Demographic", "") == "Total"
    ]
    data.sort(key=lambda x: x["value"], reverse=True)
    return {"chart": "bar", "title": "Sports Viewership — First Choice (Canada Total)", "data": data[:8], "xKey": "name", "yKey": "value"}


def _chart_nativity():
    rows = _read_csv("csfimi_sports_fandom_insights_2025.csv")
    data = []
    for r in rows:
        if r.get("Insight_Category") == "nativity_divide":
            try:
                data.append({"name": r.get("Demographic", ""), "value": float(r.get("Value", "0%").replace("%", ""))})
            except ValueError:
                pass
    return {"chart": "bar", "title": "Nativity Divide — First Choice by Origin", "data": data, "xKey": "name", "yKey": "value"}


def _chart_promotion_roi():
    data = [
        {"name": f"{r.get('Promotion_Type', '')} ({r.get('Channel', '')})", "roi": round(float(r.get("ROI", "0")), 2)}
        for r in _read_csv("promotion_roi_analysis_2025.csv")
    ]
    data.sort(key=lambda x: x["roi"], reverse=True)
    return {"chart": "bar", "title": "Promotion ROI by Type & Channel", "data": data[:10], "xKey": "name", "yKey": "roi"}


def _chart_say_do():
    data = [
        {"name": f"{r.get('Category', '')} — {r.get('Segment', '')}", "gap": round(float(r.get("Say_Do_Gap", r.get("Gap_Percentage", "0"))) * 100, 1)}
        for r in _read_csv("say_do_gap_food_beverage.csv")
    ]
    data.sort(key=lambda x: abs(x["gap"]), reverse=True)
    return {"chart": "bar", "title": "Say-Do Gap — Stated vs Actual Purchase (%)", "data": data[:10], "xKey": "name", "yKey": "gap"}


@router.get("/chart-types")
async def chart_types():
    """Available chart types for the visualize endpoint."""
    return {
        "chart_types": [
            {"id": "nps_comparison", "label": "NPS by Brand", "chart": "bar", "dataset": "brand_health_tracker_Q4_2025.csv"},
            {"id": "awareness", "label": "Brand Awareness (Aided vs Unaided)", "chart": "bar", "dataset": "brand_health_tracker_Q4_2025.csv"},
            {"id": "market_share", "label": "Market Share", "chart": "pie", "dataset": "competitive_benchmark_market_share.csv"},
            {"id": "sponsorship_roi", "label": "Sponsorship Property Scores", "chart": "bar", "dataset": "sponsorship_property_scores.csv"},
            {"id": "sports_viewership", "label": "Sports Viewership — First Choice", "chart": "bar", "dataset": "sports_viewership_crosstab_2025.csv"},
            {"id": "nativity_divide", "label": "Nativity Divide — Sports by Origin", "chart": "bar", "dataset": "csfimi_sports_fandom_insights_2025.csv"},
            {"id": "promotion_roi", "label": "Promotion ROI by Type", "chart": "bar", "dataset": "promotion_roi_analysis_2025.csv"},
            {"id": "say_do_gap", "label": "Say-Do Gap Analysis", "chart": "bar", "dataset": "say_do_gap_food_beverage.csv"},
        ]
    }


@router.post("/generate-training")
async def generate_training(request: Request):
    """Generate training pairs from seed queries for QLoRA fine-tuning."""
    body = await request.json()
    max_pairs = body.get("max_pairs", 10)

    try:
        from imi_training_generator import generate_from_seeds
        result = generate_from_seeds(max_pairs=max_pairs)
        return {"ok": True, **result}
    except Exception as e:
        logger.error(f"Training generation failed: {e}")
        return {"ok": False, "error": str(e)}


@router.get("/stats")
async def stats():
    """System health dashboard — everything at a glance."""
    import glob as _glob

    datasets = [f for f in os.listdir(IMI_DATA_DIR) if f.endswith((".csv", ".json"))] if os.path.isdir(IMI_DATA_DIR) else []

    vector_count = 0
    if os.path.exists(VECTOR_STORE_PATH):
        try:
            with open(VECTOR_STORE_PATH, "r") as f:
                store = json.loads(f.read())
                vector_count = len(store.get("documents", []))
        except Exception:
            pass

    learning_entries = 0
    if os.path.exists(LEARNING_LOG_PATH):
        with open(LEARNING_LOG_PATH, "r") as f:
            learning_entries = sum(1 for _ in f)

    training_pairs = 0
    avg_quality = 0
    if os.path.exists(TRAINING_FILE_PATH):
        qualities = []
        with open(TRAINING_FILE_PATH, "r") as f:
            for line in f:
                training_pairs += 1
                try:
                    entry = json.loads(line)
                    qualities.append(entry.get("metadata", {}).get("quality_score", 0))
                except json.JSONDecodeError:
                    pass
        avg_quality = round(sum(qualities) / len(qualities), 1) if qualities else 0

    return {
        "datasets": len(datasets),
        "dataset_files": sorted(datasets),
        "vector_documents": vector_count,
        "learning_log_entries": learning_entries,
        "training_pairs": training_pairs,
        "avg_training_quality": f"{avg_quality}/4",
        "model": "klaus-imi (Qwen 32B)",
        "embedding_model": "nomic-embed-text",
        "pipeline_stages": ["classify", "retrieve", "analyze", "synthesize", "learn"],
        "chart_types": 8,
        "endpoints": [
            "POST /klaus/imi/chat",
            "POST /klaus/imi/report",
            "POST /klaus/imi/visualize",
            "GET  /klaus/imi/chart-types",
            "GET  /klaus/imi/dashboard",
            "GET  /klaus/imi/stats",
            "POST /klaus/imi/ingest",
            "POST /klaus/imi/generate-training",
        ],
    }
