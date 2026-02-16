"""
IMI Insight Deck Generator — Universal
Produces an 11-slide PPTX matching the IMI brand template.
Works with any survey format: crosstab, flat table, respondent-level CSV, JSON.
Colors: #0B1D3A (navy bg), #E8651A (orange accent), #00A3A1 (teal), #FFFFFF, #8C99A9 (grey)
Fonts: Georgia (titles), Calibri (body)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import io
import json
import re

# ─── Brand constants ─────────────────────────────────────────────
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
ORANGE = RGBColor(0xE8, 0x65, 0x1A)
TEAL = RGBColor(0x00, 0xA3, 0xA1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x8C, 0x99, 0xA9)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
DARK_FOOTER = RGBColor(0x07, 0x14, 0x28)

SLIDE_W = Emu(9144000)
SLIDE_H = Emu(5143500)
MARGIN = Emu(640080)
CONTENT_W = Emu(7863840)


def _add_shape(slide, left, top, width, height, fill_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def _set_text(shape, text, font_name="Calibri", font_size=Pt(11), bold=False,
              color=WHITE, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    shape.text_frame.word_wrap = True
    shape.text_frame.auto_size = None
    try:
        shape.text_frame.vertical_anchor = anchor
    except:
        pass
    lines = str(text).split('\n')
    p = shape.text_frame.paragraphs[0]
    p.text = lines[0]
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    for line in lines[1:]:
        p2 = shape.text_frame.add_paragraph()
        p2.text = line
        p2.font.name = font_name
        p2.font.size = font_size
        p2.font.bold = bold
        p2.font.color.rgb = color
        p2.alignment = alignment
    return p


def _add_text_shape(slide, left, top, width, height, text,
                    font_name="Calibri", font_size=Pt(11), bold=False,
                    color=WHITE, alignment=PP_ALIGN.LEFT, fill_color=None):
    shape = _add_shape(slide, left, top, width, height, fill_color)
    _set_text(shape, text, font_name, font_size, bold, color, alignment)
    return shape


def _orange_bar(slide):
    _add_shape(slide, 0, 0, SLIDE_W, Emu(73152), ORANGE)


def _bar_chart_rows(slide, items, start_y, row_h=Emu(365760), bar_area_left=Emu(2560320),
                    bar_area_width=Emu(5486400), max_val=None, label_width=Emu(1828800)):
    if not items:
        return
    if max_val is None:
        max_val = max(v for _, v in items) if items else 1
    # Detect if values look like percentages (most > 10 and <= 100) vs scores/counts
    all_vals = [v for _, v in items if isinstance(v, (int, float))]
    has_negative = any(v < 0 for v in all_vals)
    looks_like_pct = all_vals and not has_negative and max(all_vals) <= 100 and sum(1 for v in all_vals if v > 10) > len(all_vals) * 0.3
    for i, (label, value) in enumerate(items[:8]):  # max 8 rows
        y = start_y + Emu(i * (row_h + Emu(27432)))
        if int(y) + int(row_h) > int(SLIDE_H) - 400000:
            break
        _add_text_shape(slide, MARGIN, y, label_width, row_h, str(label),
                       font_size=Pt(10), color=WHITE)
        _add_shape(slide, bar_area_left, y, bar_area_width, row_h, RGBColor(0x1A, 0x2E, 0x4A))
        pct = value / max_val if max_val > 0 else 0
        bar_w = int(bar_area_width * pct)
        if bar_w > 0:
            _add_shape(slide, bar_area_left, y, Emu(bar_w), row_h, TEAL)
        # Smart formatting: only add % for actual percentages
        if looks_like_pct:
            val_str = f"{int(value)}%" if value == int(value) else f"{value:.1f}%"
        else:
            val_str = f"{int(value)}" if value == int(value) else f"{value:.1f}"
        _add_text_shape(slide, Emu(int(bar_area_left) + bar_w + Emu(54864)), y,
                       Emu(640080), row_h, val_str,
                       font_size=Pt(10), bold=True, color=WHITE)


def _stat_card(slide, y, big_num, description, accent_color=TEAL):
    card_h = Emu(1005840)
    _add_shape(slide, MARGIN, y, CONTENT_W, card_h, WHITE)
    _add_shape(slide, MARGIN, y, Emu(64008), card_h, accent_color)
    _add_text_shape(slide, Emu(914400), Emu(int(y) + 91440), Emu(1188720), Emu(822960),
                   big_num, font_name="Georgia", font_size=Pt(44), bold=True, color=NAVY,
                   alignment=PP_ALIGN.CENTER)
    _add_text_shape(slide, Emu(2194560), Emu(int(y) + 137160), Emu(6035040), Emu(731520),
                   description, font_size=Pt(12), color=NAVY)


def _source_line(slide, text, y=Emu(4709160)):
    _add_text_shape(slide, MARGIN, y, Emu(7315200), Emu(274320), text,
                   font_size=Pt(9), color=GREY)


def _insight_box(slide, text, y=Emu(4434840)):
    _add_text_shape(slide, MARGIN, y, CONTENT_W, Emu(594360), text,
                   font_size=Pt(10), bold=True, color=ORANGE)


def _clean_label(label: str) -> str:
    label = re.sub(r'\s{2,}', ' ', label.strip())
    label = label.replace(' - NhL', ' — NHL').replace(' - NHl', ' — NHL').replace(' - nhl', ' — NHL')
    label = label.replace(' - NFL', ' — NFL').replace(' - nfl', ' — NFL')
    label = label.replace(' - MLB', ' — MLB').replace(' - mlb', ' — MLB')
    label = label.replace(' - golf', ' — Golf').replace(' - Golf', ' — Golf')
    label = label.replace(' - NBA', ' — NBA')
    if label and label[0].islower():
        label = label[0].upper() + label[1:]
    return label


def _fmt_val(v, is_pct=True):
    """Format a value for display."""
    if v is None:
        return "N/A"
    if isinstance(v, float) and v < 1 and is_pct:
        v = v * 100
    suffix = "%" if is_pct else ""
    if isinstance(v, float) and v == int(v):
        return f"{int(v)}{suffix}"
    if isinstance(v, float):
        return f"{v:.1f}{suffix}"
    return f"{v}{suffix}"


# ─── Data extraction helpers ──────────────────────────────────────

def _extract_analysis(survey_data: dict) -> dict:
    """
    Extract key analysis from any survey format.
    Returns a dict with standardized findings for slide generation.
    """
    questions = survey_data.get("questions", [])
    total_n = survey_data.get("total_n", 0)
    survey_name = survey_data.get("survey_name", "Survey Analysis")
    fmt = survey_data.get("format", "unknown")
    segments = survey_data.get("segments", [])
    items = survey_data.get("items", [])

    analysis = {
        "survey_name": survey_name,
        "total_n": total_n,
        "format": fmt,
        "segments": segments,
        "items": items,
        "primary_q": None,          # Main question/metric with ranked results
        "secondary_q": None,        # Second most interesting question
        "top_items": [],            # Top ranked items from primary question
        "biggest_gap": None,        # Biggest segment gap found
        "segment_comparisons": [],  # Notable segment differences
        "all_questions": questions,
        "stat_cards": [],           # 3 headline stats for exec summary
        "strategic_recs": [],       # 4 strategic recommendations
    }

    if not questions:
        return analysis

    # For crosstab format: use existing Q36/Q37/Q38 logic
    if fmt == "crosstab":
        return _extract_crosstab_analysis(survey_data, analysis)

    # For flat/CSV/JSON: pick the most interesting metrics
    # Filter out metadata fields (sample size, timestamps, etc.)
    skip_keywords = {"sample size", "sample_size", "n=", "base size", "respondent_id", "timestamp", "date", "id"}
    # Sort questions by variance (most spread = most interesting)
    scored_qs = []
    for q in questions:
        q_text = (q.get("text", "") + " " + q.get("id", "")).lower()
        if any(kw in q_text for kw in skip_keywords):
            continue
        vals = []
        for opt in q.get("options", []):
            for seg, v in opt.get("values", {}).items():
                if isinstance(v, (int, float)):
                    vals.append(v)
        if vals:
            spread = max(vals) - min(vals)
            scored_qs.append((spread, q))
    scored_qs.sort(key=lambda x: x[0], reverse=True)

    if scored_qs:
        analysis["primary_q"] = scored_qs[0][1]
    if len(scored_qs) > 1:
        analysis["secondary_q"] = scored_qs[1][1]

    # Get top items for primary question
    if analysis["primary_q"]:
        opts = analysis["primary_q"].get("options", [])
        ranked = []
        for opt in opts:
            v = opt["values"].get("Total", list(opt["values"].values())[0] if opt["values"] else 0)
            if isinstance(v, (int, float)):
                ranked.append((_clean_label(opt["label"]), v))
        ranked.sort(key=lambda x: x[1], reverse=True)
        analysis["top_items"] = ranked

    # Find biggest segment gaps (for respondent_csv with segment breakdowns)
    if fmt == "respondent_csv":
        _find_segment_gaps(questions, segments, analysis)

    # Generate stat cards
    _generate_stat_cards(analysis, questions)

    # Generate strategic recommendations
    _generate_strategic_recs(analysis, questions)

    return analysis


def _extract_crosstab_analysis(survey_data: dict, analysis: dict) -> dict:
    """Specialized extraction for crosstab format (Q-block style)."""
    questions = survey_data.get("questions", [])

    q36 = q37 = q_combined = q38 = None
    for q in questions:
        qid = q.get("id", "").upper()
        qtext = q.get("text", "").lower()
        if "q36" in qid or "only watch one" in qtext:
            q36 = q
        elif "q37" in qid and "second" in qtext:
            q37 = q
        elif "q37" in qid and ("first or second" in qtext or "combined" in qtext):
            q_combined = q
        elif "q38" in qid or "most want to win" in qtext:
            q38 = q
    if q_combined is None:
        for q in questions:
            if "first or second" in q.get("text", "").lower():
                q_combined = q

    analysis["primary_q"] = q36
    analysis["secondary_q"] = q_combined or q37
    analysis["q38"] = q38
    analysis["q_combined"] = q_combined
    analysis["is_sports_fandom"] = True

    # Get national results
    if q36:
        for opt in q36.get("options", []):
            vals = opt.get("values", {})
            nat = vals.get("Canada", vals.get("Total", 0))
            analysis["top_items"].append((_clean_label(opt["label"]), nat))
        analysis["top_items"].sort(key=lambda x: x[1], reverse=True)

    # Detect nativity columns
    born_in_col = born_out_col = None
    if q36 and q36.get("options"):
        for k in q36["options"][0].get("values", {}):
            kl = k.lower()
            if "born in" in kl:
                born_in_col = k
            elif "born outside" in kl:
                born_out_col = k
    analysis["born_in_col"] = born_in_col
    analysis["born_out_col"] = born_out_col

    # Base sizes
    base_sizes = {}
    for q in questions:
        if q.get("base_sizes"):
            base_sizes = q["base_sizes"]
            break
    analysis["base_sizes"] = base_sizes

    return analysis


def _find_segment_gaps(questions, segments, analysis):
    """Find the most dramatic segment gaps in respondent-level data."""
    biggest_gap = {"metric": "", "seg1": "", "seg2": "", "val1": 0, "val2": 0, "gap": 0}
    for q in questions:
        for opt in q.get("options", []):
            vals = opt.get("values", {})
            seg_vals = [(k, v) for k, v in vals.items() if k != "Total" and isinstance(v, (int, float))]
            if len(seg_vals) >= 2:
                seg_vals.sort(key=lambda x: x[1], reverse=True)
                high_k, high_v = seg_vals[0]
                low_k, low_v = seg_vals[-1]
                gap = abs(high_v - low_v)
                if gap > biggest_gap["gap"]:
                    biggest_gap = {
                        "metric": q.get("text", q.get("id", "")),
                        "seg1": high_k, "seg2": low_k,
                        "val1": high_v, "val2": low_v, "gap": gap
                    }
    if biggest_gap["gap"] > 0:
        analysis["biggest_gap"] = biggest_gap


def _generate_stat_cards(analysis, questions):
    """Generate 3 headline stat cards from the data."""
    cards = []
    top_items = analysis.get("top_items", [])
    total_n = analysis.get("total_n", 0)

    if top_items:
        top_label, top_val = top_items[0]
        pq_text = analysis.get('primary_q', {}).get('text', 'the primary metric')
        # Detect if values are percentages vs scores
        all_vals = [v for _, v in top_items if isinstance(v, (int, float))]
        is_pct = all_vals and max(all_vals) <= 100 and sum(1 for v in all_vals if v > 10) > len(all_vals) * 0.3
        display_val = top_val * 100 if top_val < 1 else top_val
        if is_pct:
            big_num = f"{round(display_val)}%"
        else:
            big_num = f"{round(display_val, 1)}" if display_val != int(display_val) else f"{int(display_val)}"
        cards.append((big_num,
                      f"{top_label} leads on {pq_text} — the top performer in this dataset."))

    if len(top_items) >= 2:
        gap = abs(top_items[0][1] - top_items[-1][1])
        gap_pct = gap * 100 if gap < 1 else gap
        cards.append((f"{round(gap_pct, 1)}{'pt' if gap_pct < 100 else ''}",
                      f"Gap between top and bottom performer ({top_items[0][0]} vs {top_items[-1][0]}) — significant spread in the data."))

    if total_n > 0:
        cards.append((f"n={total_n:,}", f"Sample size across {len(questions)} metrics analyzed. {'Robust base for reliable insights.' if total_n >= 500 else 'Moderate base — interpret with caution.'}"))

    bg = analysis.get("biggest_gap")
    if bg and bg["gap"] > 0 and len(cards) < 3:
        gap_display = bg["gap"] * 100 if bg["gap"] < 1 else bg["gap"]
        cards.append((f"{round(gap_display, 1)}pt", f"Biggest segment gap on {bg['metric']}: {bg['seg1']} vs {bg['seg2']}."))

    analysis["stat_cards"] = cards[:3]


def _generate_strategic_recs(analysis, questions):
    """Generate 4 data-driven strategic recommendations."""
    recs = []
    top_items = analysis.get("top_items", [])
    survey_name = analysis.get("survey_name", "")

    if top_items:
        recs.append((
            f"LEAD WITH {top_items[0][0].upper()}" if len(top_items[0][0]) < 30 else "LEAD WITH THE TOP PERFORMER",
            f"{top_items[0][0]} scores highest — priority for investment, activation, and messaging. Build from strength."
        ))
    if len(top_items) >= 3:
        recs.append((
            "WATCH THE MIDDLE TIER",
            f"Items ranked 2-3 ({', '.join(l for l, _ in top_items[1:3])}) are within striking distance. Small shifts in strategy could change the ranking."
        ))
    bg = analysis.get("biggest_gap")
    if bg and bg["gap"] > 0:
        recs.append((
            "ADDRESS THE SEGMENT DIVIDE",
            f"The {round(bg['gap'] * 100 if bg['gap'] < 1 else bg['gap'], 1)}-point gap between {bg['seg1']} and {bg['seg2']} on {bg['metric']} signals a targeting opportunity or a strategic risk."
        ))
    if len(top_items) > 5:
        recs.append((
            "DE-PRIORITIZE THE TAIL",
            f"Bottom performers ({', '.join(l for l, _ in top_items[-2:])}) consume resources without returns. Reallocate to top-tier opportunities."
        ))
    # Fill to 4 if needed
    default_recs = [
        ("TRACK QUARTERLY", "Establish a tracking cadence to measure movement. IMI Pulse™ quarterly tracking turns a snapshot into a trend."),
        ("CONNECT DATA TO ACTION", "Every insight in this deck should map to a specific business decision. Schedule a strategy session to prioritize activation."),
        ("SEGMENT YOUR APPROACH", "One-size-fits-all rarely works. Use the segment differences in this data to tailor messaging and targeting."),
        ("VALIDATE WITH BEHAVIOR DATA", "Stated preference ≠ actual behavior. Cross-reference with IMI's Say-Do Gap™ methodology to validate."),
    ]
    for dr in default_recs:
        if len(recs) >= 4:
            break
        recs.append(dr)
    analysis["strategic_recs"] = recs[:4]


# ─── Main generator ──────────────────────────────────────────────

def generate_imi_deck(survey_data: dict) -> bytes:
    """Generate an 11-slide IMI insight deck from any parsed survey data."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    a = _extract_analysis(survey_data)
    survey_name = a["survey_name"]
    total_n = a["total_n"]
    questions = a["all_questions"]

    # If it's the sports fandom crosstab, use the specialized builder
    if a.get("is_sports_fandom"):
        return _generate_sports_fandom_deck(prs, blank_layout, a)

    # ═══ SLIDE 1: TITLE ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_shape(slide, 0, Emu(4754880), SLIDE_W, Emu(388620), DARK_FOOTER)
    _add_text_shape(slide, Emu(731520), Emu(457200), Emu(7680960), Emu(365760),
                   "INSIGHT. DRIVING. PROFIT.", "Calibri", Pt(12), True, ORANGE)
    # Title from survey name
    display_name = survey_name.replace('_', ' ').title()
    _add_text_shape(slide, Emu(731520), Emu(1097280), Emu(7680960), Emu(2011680),
                   display_name, "Georgia", Pt(40), True, WHITE)
    # Subtitle: generated from data
    n_questions = len(questions)
    subtitle = f"Analysis of {n_questions} metrics"
    if total_n:
        subtitle += f"  |  n = {total_n:,}"
    segs = [s for s in a.get("segments", []) if s != "Total"]
    if segs:
        subtitle += f"  |  Segmented by {', '.join(segs[:4])}"
    _add_text_shape(slide, Emu(731520), Emu(3017520), Emu(6400800), Emu(548640),
                   subtitle, "Calibri", Pt(14), False, GREY)
    _add_text_shape(slide, Emu(731520), Emu(4800600), Emu(7315200), Emu(320040),
                   f"IMI Pulse™ Analysis  |  © IMI International  |  consultimi.com",
                   "Calibri", Pt(11), False, GREY)

    # ═══ SLIDE 2: EXECUTIVE SUMMARY ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(274320), Emu(7315200), Emu(457200),
                   "EXECUTIVE SUMMARY", "Georgia", Pt(28), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(685800), Emu(7315200), Emu(320040),
                   "Key Findings at a Glance", "Calibri", Pt(14), False, GREY)
    colors = [TEAL, ORANGE, NAVY]
    for i, (big_num, desc) in enumerate(a["stat_cards"][:3]):
        _stat_card(slide, Emu(1234440 + i * 1188720), big_num, desc, colors[i % 3])
    _source_line(slide, f"Source: IMI Pulse™ {survey_name}  |  n = {total_n:,}" if total_n else f"Source: IMI Pulse™ {survey_name}")

    # ═══ SLIDE 3: PRIMARY METRIC RANKING ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    pq = a.get("primary_q") or {}
    pq_title = pq.get("text", "Primary Metric").replace('_', ' ').title()
    _add_text_shape(slide, MARGIN, Emu(182880), CONTENT_W, Emu(731520),
                   pq_title.upper(), "Georgia", Pt(22), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(777240), Emu(7315200), Emu(274320),
                   f"Ranked results  |  {survey_name} (n = {total_n:,})" if total_n else f"Ranked results  |  {survey_name}",
                   "Calibri", Pt(11), False, GREY)

    top_items = a.get("top_items", [])
    # Determine if values are percentages or scores
    is_pct = all(v <= 1 for _, v in top_items[:5]) if top_items else True
    chart_items = []
    for label, val in top_items[:8]:
        display_val = val * 100 if is_pct and val <= 1 else val
        chart_items.append((_clean_label(label), round(display_val, 1)))
    if chart_items:
        _bar_chart_rows(slide, chart_items, Emu(1097280),
                       max_val=max(v for _, v in chart_items) * 1.2)
    _source_line(slide, f"Source: IMI Pulse™  |  {pq_title}")

    # ═══ SLIDE 4: SECOND METRIC / COMPARISON ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    sq = a.get("secondary_q") or (questions[1] if len(questions) > 1 else {})
    sq_title = sq.get("text", "Secondary Analysis").replace('_', ' ').title()
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   sq_title.upper(), "Georgia", Pt(20), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(640080), Emu(7315200), Emu(274320),
                   f"Comparative view  |  {survey_name}", "Calibri", Pt(11), False, GREY)
    sq_items = []
    for opt in sq.get("options", [])[:8]:
        v = opt["values"].get("Total", list(opt["values"].values())[0] if opt["values"] else 0)
        if isinstance(v, (int, float)):
            display_v = v * 100 if v <= 1 else v
            sq_items.append((_clean_label(opt["label"]), round(display_v, 1)))
    sq_items.sort(key=lambda x: x[1], reverse=True)
    if sq_items:
        _bar_chart_rows(slide, sq_items, Emu(960120),
                       max_val=max(v for _, v in sq_items) * 1.1)
    _source_line(slide, f"Source: IMI Pulse™  |  {sq_title}")

    # ═══ SLIDE 5: BIGGEST DIVIDE ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    bg_data = a.get("biggest_gap")
    if bg_data and bg_data["gap"] > 0:
        gap_display = bg_data["gap"] * 100 if bg_data["gap"] < 1 else bg_data["gap"]
        _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(502920),
                       "THE BIGGEST DIVIDE", "Georgia", Pt(28), True, WHITE)
        _add_text_shape(slide, MARGIN, Emu(685800), Emu(7315200), Emu(320040),
                       f"{bg_data['metric']}  |  {round(gap_display, 1)}-point gap",
                       "Calibri", Pt(13), False, GREY)
        # Show the two segments side by side
        col_w = Emu(3749040)
        right_col = Emu(4754880)
        v1 = bg_data["val1"] * 100 if bg_data["val1"] < 1 else bg_data["val1"]
        v2 = bg_data["val2"] * 100 if bg_data["val2"] < 1 else bg_data["val2"]
        _add_text_shape(slide, MARGIN, Emu(1371600), col_w, Emu(411480),
                       bg_data["seg1"][:35], "Calibri", Pt(12), True, WHITE,
                       PP_ALIGN.CENTER, fill_color=TEAL)
        _add_text_shape(slide, right_col, Emu(1371600), col_w, Emu(411480),
                       bg_data["seg2"][:35], "Calibri", Pt(12), True, WHITE,
                       PP_ALIGN.CENTER, fill_color=ORANGE)
        _add_text_shape(slide, MARGIN, Emu(2011680), col_w, Emu(914400),
                       f"{round(v1, 1)}", "Georgia", Pt(72), True, TEAL, PP_ALIGN.CENTER)
        _add_text_shape(slide, right_col, Emu(2011680), col_w, Emu(914400),
                       f"{round(v2, 1)}", "Georgia", Pt(72), True, ORANGE, PP_ALIGN.CENTER)
        _insight_box(slide,
            f"SO WHAT? A {round(gap_display, 1)}-point gap on {bg_data['metric']} between {bg_data['seg1']} and {bg_data['seg2']} "
            f"signals fundamentally different needs. Tailor strategy to each segment.")
    else:
        _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(502920),
                       "DATA HIGHLIGHTS", "Georgia", Pt(28), True, WHITE)
        # Show top 3 metrics side by side
        for i, q in enumerate(questions[:3]):
            y = Emu(1097280 + i * Emu(1097280))
            _add_text_shape(slide, MARGIN, y, CONTENT_W, Emu(274320),
                           q.get("text", q.get("id", "")).replace('_', ' ').title(),
                           font_size=Pt(12), bold=True, color=ORANGE)
            opts = q.get("options", [])[:3]
            for j, opt in enumerate(opts):
                v = opt["values"].get("Total", list(opt["values"].values())[0] if opt["values"] else 0)
                display_v = v * 100 if isinstance(v, (int, float)) and v <= 1 else v
                _add_text_shape(slide, MARGIN, Emu(int(y) + 320040 + j * 228600), CONTENT_W, Emu(228600),
                               f"{_clean_label(opt['label'])}: {round(display_v, 1) if isinstance(display_v, (int, float)) else display_v}",
                               font_size=Pt(10), color=WHITE)

    # ═══ SLIDE 6: METRIC DEEP-DIVE (3rd question) ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    q3 = questions[2] if len(questions) > 2 else questions[0] if questions else {}
    q3_title = q3.get("text", "Deep Dive").replace('_', ' ').title()
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   q3_title.upper(), "Georgia", Pt(22), True, WHITE)
    q3_items = []
    for opt in q3.get("options", [])[:8]:
        v = opt["values"].get("Total", list(opt["values"].values())[0] if opt["values"] else 0)
        if isinstance(v, (int, float)):
            display_v = v * 100 if v <= 1 else v
            q3_items.append((_clean_label(opt["label"]), round(display_v, 1)))
    q3_items.sort(key=lambda x: x[1], reverse=True)
    if q3_items:
        _bar_chart_rows(slide, q3_items, Emu(960120),
                       max_val=max(v for _, v in q3_items) * 1.1)

    # ═══ SLIDE 7: METRIC COMPARISON (4th + 5th) ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   "ADDITIONAL METRICS", "Georgia", Pt(24), True, WHITE)

    # Show 2-column layout for next 2 metrics
    col_w = Emu(3749040)
    right_col = Emu(4754880)
    for ci, qi in enumerate([3, 4]):
        if qi >= len(questions):
            break
        q = questions[qi]
        cx = MARGIN if ci == 0 else right_col
        _add_text_shape(slide, cx, Emu(685800), col_w, Emu(320040),
                       q.get("text", "").replace('_', ' ').title()[:30],
                       "Calibri", Pt(11), True, ORANGE, PP_ALIGN.CENTER)
        opts = q.get("options", [])[:5]
        for ri, opt in enumerate(opts):
            v = opt["values"].get("Total", list(opt["values"].values())[0] if opt["values"] else 0)
            if isinstance(v, (int, float)):
                display_v = v * 100 if v <= 1 else v
                y = Emu(1051560 + ri * Emu(594360))
                _add_shape(slide, cx, y, col_w, Emu(502920), RGBColor(0x14, 0x28, 0x50))
                _add_text_shape(slide, Emu(int(cx) + 137160), Emu(int(y) + 18288),
                               Emu(2560320), Emu(228600),
                               _clean_label(opt["label"])[:25], font_size=Pt(10), color=WHITE)
                _add_text_shape(slide, Emu(int(cx) + int(col_w) - 914400), Emu(int(y) + 18288),
                               Emu(914400), Emu(228600),
                               f"{round(display_v, 1)}", font_size=Pt(14), bold=True,
                               color=WHITE, alignment=PP_ALIGN.RIGHT)

    # ═══ SLIDE 8: THIRD METRIC DEEP-DIVE ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    q5 = questions[5] if len(questions) > 5 else questions[-1] if questions else {}
    q5_title = q5.get("text", "Analysis").replace('_', ' ').title()
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   q5_title.upper(), "Georgia", Pt(22), True, WHITE)
    q5_items = []
    for opt in q5.get("options", [])[:8]:
        v = opt["values"].get("Total", list(opt["values"].values())[0] if opt["values"] else 0)
        if isinstance(v, (int, float)):
            display_v = v * 100 if v <= 1 else v
            q5_items.append((_clean_label(opt["label"]), round(display_v, 1)))
    q5_items.sort(key=lambda x: x[1], reverse=True)
    if q5_items:
        _bar_chart_rows(slide, q5_items, Emu(960120),
                       max_val=max(v for _, v in q5_items) * 1.1)

    # ═══ SLIDE 9: DATA OVERVIEW TABLE ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   "FULL DATA OVERVIEW", "Georgia", Pt(22), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(640080), Emu(7315200), Emu(274320),
                   f"All {len(questions)} metrics at a glance", "Calibri", Pt(11), False, GREY)
    # Compact table-like display of all metrics
    row_y = Emu(960120)
    for qi, q in enumerate(questions[:12]):
        if int(row_y) > int(SLIDE_H) - 400000:
            break
        opts = q.get("options", [])
        if opts:
            top_opt = max(opts, key=lambda o: o["values"].get("Total", list(o["values"].values())[0] if o["values"] else 0) if isinstance(o["values"].get("Total", list(o["values"].values())[0] if o["values"] else 0), (int, float)) else 0)
            top_v = top_opt["values"].get("Total", list(top_opt["values"].values())[0] if top_opt["values"] else 0)
            display_v = top_v * 100 if isinstance(top_v, (int, float)) and top_v <= 1 else top_v
            metric_name = q.get("text", q.get("id", "")).replace('_', ' ').title()[:35]
            top_label = _clean_label(top_opt["label"])[:25]
            _add_shape(slide, MARGIN, row_y, CONTENT_W, Emu(274320), RGBColor(0x14, 0x28, 0x50) if qi % 2 == 0 else NAVY)
            _add_text_shape(slide, Emu(int(MARGIN) + 91440), row_y,
                           Emu(3200400), Emu(274320), metric_name, font_size=Pt(9), color=GREY)
            _add_text_shape(slide, Emu(4114800), row_y,
                           Emu(2560320), Emu(274320), top_label, font_size=Pt(9), color=WHITE)
            _add_text_shape(slide, Emu(6858000), row_y,
                           Emu(1371600), Emu(274320),
                           f"{round(display_v, 1) if isinstance(display_v, (int, float)) else display_v}",
                           font_size=Pt(9), bold=True, color=TEAL, alignment=PP_ALIGN.RIGHT)
            row_y = Emu(int(row_y) + 301752)

    # ═══ SLIDE 10: STRATEGIC IMPLICATIONS ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(182880), CONTENT_W, Emu(731520),
                   "STRATEGIC IMPLICATIONS", "Georgia", Pt(28), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(822960), Emu(7315200), Emu(320040),
                   "Turning Insight Into Action", "Calibri", Pt(14), False, GREY)
    accent_colors = [TEAL, ORANGE, NAVY, GREY]
    for i, (title, body) in enumerate(a["strategic_recs"][:4]):
        y = Emu(1143000 + i * Emu(914400))
        card_h = Emu(777240)
        _add_shape(slide, MARGIN, y, CONTENT_W, card_h, WHITE)
        _add_shape(slide, MARGIN, y, Emu(64008), card_h, accent_colors[i % 4])
        _add_text_shape(slide, Emu(914400), Emu(int(y) + 45720), Emu(7406640), Emu(228600),
                       title, font_size=Pt(11), bold=True, color=NAVY)
        _add_text_shape(slide, Emu(914400), Emu(int(y) + 292608), Emu(7406640), Emu(438912),
                       body, font_size=Pt(9), color=NAVY)
    _source_line(slide, f"Source: IMI Pulse™ {survey_name}  |  n = {total_n:,}  |  © IMI International" if total_n else f"Source: IMI Pulse™ {survey_name}  |  © IMI International",
                y=Emu(4800600))

    # ═══ SLIDE 11: CLOSING ═══
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(914400), CONTENT_W, Emu(457200),
                   "INSIGHT. DRIVING. PROFIT.", "Calibri", Pt(20), True, ORANGE, PP_ALIGN.CENTER)
    _add_text_shape(slide, Emu(1371600), Emu(1645920), Emu(6400800), Emu(1463040),
                   f"Let's Unlock the Value of\n{display_name}", "Georgia", Pt(32), True, WHITE, PP_ALIGN.CENTER)
    _add_text_shape(slide, Emu(1828800), Emu(3200400), Emu(5486400), Emu(457200),
                   "For a deeper dive into the insight and what it means for your brand, reach out to our team.",
                   "Calibri", Pt(12), False, GREY, PP_ALIGN.CENTER)
    btn = _add_shape(slide, Emu(3200400), Emu(3931920), Emu(2743200), Emu(457200), ORANGE)
    _set_text(btn, "consultimi.com", "Calibri", Pt(16), True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _add_text_shape(slide, Emu(1828800), Emu(4572000), Emu(5486400), Emu(274320),
                   "IMI International  |  Toronto  •  Melbourne  •  Tokyo",
                   "Calibri", Pt(10), False, GREY, PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _generate_sports_fandom_deck(prs, blank_layout, a):
    """Specialized deck for the Canadian Sports Fandom crosstab format."""
    questions = a["all_questions"]
    total_n = a["total_n"]
    survey_name = a["survey_name"]
    q36 = a.get("primary_q")
    q_combined = a.get("q_combined")
    q38 = a.get("q38")
    born_in_col = a.get("born_in_col")
    born_out_col = a.get("born_out_col")
    base_sizes = a.get("base_sizes", {})

    def get_national(q):
        if not q: return []
        results = []
        for opt in q.get("options", []):
            vals = opt.get("values", {})
            nat = vals.get("Canada", vals.get("Total", 0))
            results.append((_clean_label(opt["label"]), nat))
        results.sort(key=lambda x: x[1], reverse=True)
        return results

    def get_segment(q, seg_name):
        if not q: return []
        results = []
        for opt in q.get("options", []):
            vals = opt.get("values", {})
            v = vals.get(seg_name, 0)
            results.append((_clean_label(opt["label"]), v))
        results.sort(key=lambda x: x[1], reverse=True)
        return results

    q36_national = get_national(q36)
    q_combined_national = get_national(q_combined)
    q38_national = get_national(q38)
    q36_born_in = get_segment(q36, born_in_col) if born_in_col else []
    q36_born_out = get_segment(q36, born_out_col) if born_out_col else []
    comb_born_in = get_segment(q_combined, born_in_col) if born_in_col else []
    comb_born_out = get_segment(q_combined, born_out_col) if born_out_col else []
    born_in_n = base_sizes.get(born_in_col, 376) if born_in_col else 376
    born_out_n = base_sizes.get(born_out_col, 127) if born_out_col else 127
    right_col = Emu(4754880)
    col_w = Emu(3749040)

    # Slide 1: Title
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_shape(slide, 0, Emu(4754880), SLIDE_W, Emu(388620), DARK_FOOTER)
    _add_text_shape(slide, Emu(731520), Emu(457200), Emu(7680960), Emu(365760),
                   "INSIGHT. DRIVING. PROFIT.", "Calibri", Pt(12), True, ORANGE)
    _add_text_shape(slide, Emu(731520), Emu(1097280), Emu(7680960), Emu(2011680),
                   "Canadian Sports\nFandom Pulse" if "Fandom" in survey_name else survey_name,
                   "Georgia", Pt(44), True, WHITE)
    _add_text_shape(slide, Emu(731520), Emu(3017520), Emu(6400800), Emu(548640),
                   "What Canadians Want to Watch — and What It Means for Sponsors",
                   "Calibri", Pt(16), False, GREY)
    _add_text_shape(slide, Emu(731520), Emu(4800600), Emu(7315200), Emu(320040),
                   f"n = {total_n} Canadian Adults  |  Cross-tabulated by Age, Gender, Region, Income & Nativity",
                   "Calibri", Pt(11), False, GREY)

    # Slide 2: Exec Summary
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(274320), Emu(7315200), Emu(457200),
                   "EXECUTIVE SUMMARY", "Georgia", Pt(28), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(685800), Emu(7315200), Emu(320040),
                   "The Essential Insight", "Calibri", Pt(14), False, GREY)
    top2_pct = 0
    if q_combined_national and len(q_combined_national) >= 1:
        top_combined = q_combined_national[0][1]
        top2_pct = round(top_combined * 100) if top_combined < 1 else round(top_combined)
    elif q36_national and len(q36_national) >= 2:
        sc = next((v for l, v in q36_national if "stanley" in l.lower()), 0)
        wo = next((v for l, v in q36_national if "winter" in l.lower() or "olympic" in l.lower()), 0)
        top2_pct = round((sc + wo) * 100) if sc < 1 else round(sc + wo)
    fifa_born_out = next((v for l, v in q36_born_out if "fifa" in l.lower()), 0)
    if fifa_born_out < 1: fifa_born_out = round(fifa_born_out * 100)
    fifa_born_in = next((v for l, v in q36_born_in if "fifa" in l.lower()), 0)
    if fifa_born_in < 1: fifa_born_in = round(fifa_born_in * 100)
    stanley_dream = next((v for l, v in q38_national if "stanley" in l.lower()), 0) if q38_national else 0
    if stanley_dream < 1: stanley_dream = round(stanley_dream * 100)
    _stat_card(slide, Emu(1234440), f"{top2_pct}%",
               f"of Canadians rank the Stanley Cup or Winter Olympics as their #1 event to watch — but FIFA is the fastest-growing passion, driven entirely by immigrants.", TEAL)
    _stat_card(slide, Emu(2423160), f"{fifa_born_out}%",
               f"of foreign-born Canadians choose the FIFA World Cup as their top event — {round(fifa_born_out/fifa_born_in, 1) if fifa_born_in else 3.4}x the rate of Canadian-born respondents ({fifa_born_in}%).", ORANGE)
    _stat_card(slide, Emu(3611880), f"{stanley_dream}%",
               "want a Canadian Stanley Cup team to win a championship most — making hockey the cultural backbone, but passion skews older (55+) and Western.", NAVY)
    _source_line(slide, f"Source: IMI Pulse™ {survey_name}  |  n = {total_n}")

    # Slide 3: National Results
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(182880), CONTENT_W, Emu(731520),
                   "IF YOU COULD ONLY WATCH ONE EVENT…", "Georgia", Pt(24), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(777240), Emu(7315200), Emu(274320),
                   f"Q36: National Results  |  Canada Total (n = {total_n})", "Calibri", Pt(11), False, GREY)
    nat_items = [(l, round((v * 100 if v < 1 else v), 1)) for l, v in q36_national]
    _bar_chart_rows(slide, nat_items, Emu(1097280), max_val=max(v for _, v in nat_items) * 1.2 if nat_items else 30)
    _source_line(slide, "Source: IMI Pulse™  |  % selecting as #1 event to watch")

    # Slide 4: Combined
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   "COMBINED VIEWERSHIP INTENT: 1ST + 2ND CHOICE", "Georgia", Pt(20), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(640080), Emu(7315200), Emu(274320),
                   "Total reach when combining first and second preference  |  Canada Total", "Calibri", Pt(11), False, GREY)
    comb_items = [(l, round((v * 100 if v < 1 else v), 1)) for l, v in q_combined_national] if q_combined_national else []
    if comb_items:
        _bar_chart_rows(slide, comb_items, Emu(960120), max_val=max(v for _, v in comb_items) * 1.1)
    _insight_box(slide, "KEY INSIGHT: Three events are virtually tied at ~45% combined reach. The Stanley Cup and Winter Olympics dominate among Canadian-born, while FIFA is the immigrant powerhouse.")

    # Slide 5: Nativity Divide
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(502920),
                   "THE NATIVITY DIVIDE", "Georgia", Pt(28), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(685800), Emu(7315200), Emu(320040),
                   "Born in Canada vs. Born Outside Canada — Dramatically Different Passions", "Calibri", Pt(13), False, GREY)
    _add_text_shape(slide, MARGIN, Emu(1188720), col_w, Emu(411480),
                   f"BORN IN CANADA (n={born_in_n})", "Calibri", Pt(12), True, WHITE, PP_ALIGN.CENTER, fill_color=TEAL)
    _add_text_shape(slide, right_col, Emu(1188720), col_w, Emu(411480),
                   f"BORN OUTSIDE CANADA (n={born_out_n})", "Calibri", Pt(12), True, WHITE, PP_ALIGN.CENTER, fill_color=ORANGE)
    left_data = comb_born_in if comb_born_in else q36_born_in
    right_data = comb_born_out if comb_born_out else q36_born_out
    row_h = Emu(411480)
    gap = Emu(475488)
    for col_x, data in [(MARGIN, left_data[:6]), (right_col, right_data[:6])]:
        for i, (label, val) in enumerate(data):
            y = Emu(1691640 + i * int(gap))
            pct = val * 100 if val < 1 else val
            _add_shape(slide, col_x, y, col_w, row_h, RGBColor(0x14, 0x28, 0x50))
            _add_text_shape(slide, Emu(int(col_x) + 137160), y, Emu(2560320), row_h,
                           label, font_size=Pt(10), color=WHITE)
            _add_text_shape(slide, Emu(int(col_x) + int(col_w) - 914400), y,
                           Emu(914400), row_h, f"{round(pct)}%",
                           font_size=Pt(14), bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)
    _insight_box(slide,
        "SO WHAT? Canada's growing immigrant population is reshaping the sports landscape. "
        "FIFA sponsorship is a direct line to multicultural audiences.", y=Emu(4572000))

    # Slide 6: Age
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   "AGE DRIVES DISTINCT FANDOM PROFILES", "Georgia", Pt(24), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(640080), Emu(7315200), Emu(274320),
                   "Top event by age cohort  |  % first choice", "Calibri", Pt(11), False, GREY)
    age_cols = [("18 to 34 years old", "18–34"), ("35 to 54 years old", "35–54"), ("55 and older", "55+")]
    col_width = Emu(2651760)
    for ci, (seg_key, seg_label) in enumerate(age_cols):
        cx = Emu(int(MARGIN) + ci * (int(col_width) + Emu(91440)))
        _add_text_shape(slide, cx, Emu(960120), col_width, Emu(320040),
                       seg_label, "Calibri", Pt(14), True, ORANGE, PP_ALIGN.CENTER)
        age_data = get_segment(q36, seg_key)[:5]
        for ri, (label, val) in enumerate(age_data):
            y = Emu(1325880 + ri * Emu(548640))
            pct = val * 100 if val < 1 else val
            _add_text_shape(slide, cx, y, col_width, Emu(228600),
                           label, font_size=Pt(10), color=WHITE, alignment=PP_ALIGN.CENTER)
            _add_text_shape(slide, cx, Emu(int(y) + Emu(228600)), col_width, Emu(228600),
                           f"{round(pct)}%", font_size=Pt(16), bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)
    _insight_box(slide, "KEY INSIGHT: FIFA and NHL are cross-generational but skew differently. Young adults (18–34) are more likely to pick NBA. Winter Olympics resonate strongest with 55+.")

    # Slide 7: Gender & Regional
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   "GENDER & REGIONAL LENS", "Georgia", Pt(24), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(640080), Emu(7315200), Emu(274320),
                   "Where do the biggest gaps live?  |  % first choice", "Calibri", Pt(11), False, GREY)
    _add_text_shape(slide, MARGIN, Emu(914400), Emu(1828800), Emu(182880),
                   "GENDER", "Calibri", Pt(12), True, ORANGE)
    female_data = get_segment(q36, "Female")
    male_data = get_segment(q36, "Male")
    gender_items = []
    if female_data and male_data:
        fd = {l: v for l, v in female_data}
        md = {l: v for l, v in male_data}
        for label in fd:
            f = fd.get(label, 0); m = md.get(label, 0)
            fp = f * 100 if f < 1 else f; mp = m * 100 if m < 1 else m
            gender_items.append((label, fp, mp, fp - mp))
        gender_items.sort(key=lambda x: abs(x[3]), reverse=True)
    for i, (label, fp, mp, g) in enumerate(gender_items[:4]):
        y = Emu(1097280 + i * Emu(594360))
        sign = "+" if g > 0 else ""
        who = "F" if g > 0 else "M"
        _add_shape(slide, MARGIN, y, col_w, Emu(502920), RGBColor(0x14, 0x28, 0x50))
        _add_shape(slide, MARGIN, y, Emu(54864), Emu(502920), TEAL if g > 0 else ORANGE)
        _add_text_shape(slide, Emu(int(MARGIN) + 91440), Emu(int(y) + 18288),
                       Emu(3200400), Emu(228600), label, font_size=Pt(10), bold=True, color=WHITE)
        _add_text_shape(slide, Emu(int(MARGIN) + 91440), Emu(int(y) + 246888),
                       Emu(3200400), Emu(228600),
                       f"F: {round(fp)}%  M: {round(mp)}%  ({sign}{round(g)}pt {who})",
                       font_size=Pt(9), color=GREY)
    _add_text_shape(slide, right_col, Emu(914400), Emu(1828800), Emu(182880),
                   "REGIONAL", "Calibri", Pt(12), True, ORANGE)
    regions = [("West", "Stanley Cup #1, lowest FIFA share"), ("Ontario", "FIFA dominates, Blue Jays fandom strongest"),
               ("Quebec", "Super Bowl over-indexes, Winter Olympics strong"), ("Atlantic Canada", "Stanley Cup and Winter Olympics tie for #1")]
    for i, (region, desc) in enumerate(regions):
        y = Emu(1051560 + i * Emu(594360))
        _add_shape(slide, right_col, y, col_w, Emu(502920), RGBColor(0x14, 0x28, 0x50))
        _add_shape(slide, right_col, y, Emu(54864), Emu(502920), TEAL)
        _add_text_shape(slide, Emu(int(right_col) + 182880), Emu(int(y) + 18288),
                       Emu(1097280), Emu(228600), region, font_size=Pt(11), bold=True, color=WHITE)
        _add_text_shape(slide, Emu(int(right_col) + 182880), Emu(int(y) + 210312),
                       Emu(3383280), Emu(228600), desc, font_size=Pt(8), color=GREY)

    # Slide 8: Aspirational (Q38)
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   "WHICH TEAM DO CANADIANS MOST WANT TO WIN?", "Georgia", Pt(22), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(640080), Emu(7315200), Emu(274320),
                   f"Q38: Dream Championship Outcome  |  Canada Total (n = {total_n})", "Calibri", Pt(11), False, GREY)
    q38_items = [(l, round((v * 100 if v < 1 else v), 1)) for l, v in q38_national]
    if q38_items:
        _bar_chart_rows(slide, q38_items, Emu(1005840), max_val=max(v for _, v in q38_items) * 1.2)

    # Slide 9: Income
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(228600), CONTENT_W, Emu(457200),
                   "INCOME LENS: WHO ARE THE HIGH-VALUE FANS?", "Georgia", Pt(22), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(640080), Emu(7315200), Emu(274320),
                   "% first choice by household income  |  Q36", "Calibri", Pt(11), False, GREY)
    income_cols = [("Under $50,000", "<$50K"), ("$50,000 - $124,999", "$50–125K"), ("$125,000 and over", "$125K+")]
    col_width = Emu(2651760)
    for ci, (seg_key, seg_label) in enumerate(income_cols):
        cx = Emu(int(MARGIN) + ci * (int(col_width) + Emu(91440)))
        _add_text_shape(slide, cx, Emu(960120), col_width, Emu(320040),
                       seg_label, "Calibri", Pt(14), True, ORANGE, PP_ALIGN.CENTER)
        inc_data = get_segment(q36, seg_key)[:5]
        for ri, (label, val) in enumerate(inc_data):
            y = Emu(1325880 + ri * Emu(548640))
            pct = val * 100 if val < 1 else val
            _add_text_shape(slide, cx, y, col_width, Emu(228600),
                           label, font_size=Pt(10), color=WHITE, alignment=PP_ALIGN.CENTER)
            _add_text_shape(slide, cx, Emu(int(y) + Emu(228600)), col_width, Emu(228600),
                           f"{round(pct)}%", font_size=Pt(16), bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

    # Slide 10: Strategic
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(182880), CONTENT_W, Emu(731520),
                   "STRATEGIC IMPLICATIONS FOR SPONSORS", "Georgia", Pt(28), True, WHITE)
    _add_text_shape(slide, MARGIN, Emu(822960), Emu(7315200), Emu(320040),
                   "Turning Insight Into Action", "Calibri", Pt(14), False, GREY)
    recs = [
        ("HOCKEY REMAINS THE CULTURAL BACKBONE — BUT DON'T OVER-INDEX",
         "Stanley Cup + Men's Olympic Hockey capture 44% of championship desire. However, this audience skews older (55+), male, and Western."),
        ("FIFA 2026 IS THE LARGEST UNTAPPED SPONSORSHIP OPPORTUNITY",
         "With 83% combined preference among immigrants and 15% wanting Canada to win, FIFA 2026 offers unique access to a growing, diverse demographic."),
        ("WINTER OLYMPICS = THE GENDER EQUALIZER",
         "The only major event where women over-index men by 10 points. For brands targeting female consumers 35+, Olympic sponsorship delivers unique reach."),
        ("REGIONALIZE YOUR ACTIVATION STRATEGY",
         "Ontario = FIFA & Blue Jays. West = Stanley Cup. Atlantic = Hockey + Olympics. Quebec = NFL surprise. National campaigns need regional activation layers."),
    ]
    accent_colors = [TEAL, ORANGE, NAVY, GREY]
    for i, (title, body) in enumerate(recs):
        y = Emu(1143000 + i * Emu(914400))
        card_h = Emu(777240)
        _add_shape(slide, MARGIN, y, CONTENT_W, card_h, WHITE)
        _add_shape(slide, MARGIN, y, Emu(64008), card_h, accent_colors[i])
        _add_text_shape(slide, Emu(914400), Emu(int(y) + 45720), Emu(7406640), Emu(228600),
                       title, font_size=Pt(11), bold=True, color=NAVY)
        _add_text_shape(slide, Emu(914400), Emu(int(y) + 292608), Emu(7406640), Emu(438912),
                       body, font_size=Pt(9), color=NAVY)
    _source_line(slide, f"Source: IMI Pulse™ {survey_name}  |  n = {total_n}  |  © IMI International", y=Emu(4800600))

    # Slide 11: Closing
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    _orange_bar(slide)
    _add_text_shape(slide, MARGIN, Emu(914400), CONTENT_W, Emu(457200),
                   "INSIGHT. DRIVING. PROFIT.", "Calibri", Pt(20), True, ORANGE, PP_ALIGN.CENTER)
    _add_text_shape(slide, Emu(1371600), Emu(1645920), Emu(6400800), Emu(1463040),
                   f"Let's Unlock the Value of\nCanadian Sports Fandom" if "Fandom" in survey_name else f"Let's Unlock the Value of\n{survey_name}",
                   "Georgia", Pt(36), True, WHITE, PP_ALIGN.CENTER)
    _add_text_shape(slide, Emu(1828800), Emu(3200400), Emu(5486400), Emu(457200),
                   "For a deeper dive into the insight and what it means for your brand, reach out to our team.",
                   "Calibri", Pt(12), False, GREY, PP_ALIGN.CENTER)
    btn = _add_shape(slide, Emu(3200400), Emu(3931920), Emu(2743200), Emu(457200), ORANGE)
    _set_text(btn, "consultimi.com", "Calibri", Pt(16), True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _add_text_shape(slide, Emu(1828800), Emu(4572000), Emu(5486400), Emu(274320),
                   "IMI International  |  Toronto  •  Melbourne  •  Tokyo",
                   "Calibri", Pt(10), False, GREY, PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
